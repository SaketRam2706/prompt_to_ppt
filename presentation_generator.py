import json
import requests
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from io import BytesIO
from openai import OpenAI
import os  # <-- Add this
from dotenv import load_dotenv
load_dotenv()

# --- API Keys are now loaded from environment variables ---
# Set DEEPSEEK_API_KEY and UNSPLASH_ACCESS_KEY in your environment or in a .env file
# Example .env:
# DEEPSEEK_API_KEY=sk-xxxxxxx
# UNSPLASH_ACCESS_KEY=xxxxxxx

UNSPLASH_ACCESS_KEY = os.environ.get("UNSPLASH_ACCESS_KEY")
DEEPSEEK_API_KEY = os.environ.get("DEEPSEEK_API_KEY")

# ────── DeepSeek Presentation Content Generation ──────
def get_presentation_content(prompt: str, min_slides=5, max_slides=10) -> dict:
    """
    Generates presentation content using the DeepSeek API.
    """
    client = OpenAI(api_key=DEEPSEEK_API_KEY, base_url="https://api.deepseek.com")
    system_prompt = f"""
    You are an expert in creating presentations.
    For each slide, include:
    - title
    - content (a list of bullet points)
    - image_query for API call to Unsplash
      - Use a single string for regular slides
      - Use a list of strings for 3_BLOCK_LAYOUT or 4_BLOCK_LAYOUT
    - layout (choose from: 1ST_SLIDE, BASIC_CONTENT_SLIDE, 2/3RD_IMAGE_LEFT, 2/3RD_IMAGE_RIGHT, HALF_IMAGE_LEFT, HALF_IMAGE_RIGHT, 3_BLOCK_LAYOUT, 4_BLOCK_LAYOUT)
    Return JSON like this:
    {{
      \"title\": \"Presentation Title\",
      \"slides\": [
        {{
          \"title\": \"Slide 1 Title\",
          \"content\": [\"Point A\", \"Point B\"],
          \"image_query\": \"A futuristic city skyline at sunset\" 
          // or if layout is 3_BLOCK_LAYOUT or 4_BLOCK_LAYOUT:
          // [\"query1\", \"query2\", \"query3\"] or [\"query1\", \"query2\", \"query3\", \"query4\"] 
          \"layout\": \"2/3RD_IMAGE_RIGHT\"
        }}
      ]
    }}
    Limit the presentation to between {min_slides} and {max_slides} slides.
    """
    response = client.chat.completions.create(
        model="deepseek-chat",
        messages=[
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": prompt},
        ],
        max_tokens=2048,
        temperature=0.7,
        response_format={"type": "json_object"},
    )
    try:
        content = response.choices[0].message.content
        return json.loads(content)
    except (json.JSONDecodeError, KeyError) as e:
        print(f"Error parsing JSON from API response: {e}")
        print(f"Raw response: {content}")
        return None

# ────── Unsplash Image Search ──────
def get_unsplash_image_url(query: str) -> str:
    url = f"https://api.unsplash.com/search/photos?client_id={UNSPLASH_ACCESS_KEY}&query={query}"
    response = requests.get(url)
    data = response.json()
    if data.get('results'):
        return data['results'][0]['urls']['raw'] + "&w=1600&dpr=2"
    return None

def get_unsplash_image_urls(queries):
    if isinstance(queries, str):
        return get_unsplash_image_url(queries)
    elif isinstance(queries, list):
        return [get_unsplash_image_url(q) for q in queries]
    return None

# ────── Slide Layout Functions ──────
def add_image_text_slide(prs, image_url, title, body, image_pos='right', image_fraction=2/3, font_name="Arial"):
    slide_width = prs.slide_width.inches
    slide_height = prs.slide_height.inches
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    image_width = slide_width * image_fraction
    text_width = slide_width - image_width
    if image_pos == 'right':
        img_x = slide_width - image_width
        txt_x = 0
    else:
        img_x = 0
        txt_x = image_width
    # Validate image_url
    if not image_url or not isinstance(image_url, str) or not image_url.startswith("http"):
        image_url = "https://via.placeholder.com/300x200?text=No+Image"
    img_data = BytesIO(requests.get(image_url, verify=False).content)
    slide.shapes.add_picture(
        img_data,
        Inches(img_x),
        Inches(0),
        Inches(image_width),
        Inches(slide_height)
    )
    textbox = slide.shapes.add_textbox(
        Inches(txt_x),
        Inches(0),
        Inches(text_width),
        Inches(slide_height)
    )
    tf = textbox.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = title
    p.font.bold = True
    p.font.size = Pt(32)
    p.font.name = font_name
    p.font.color.rgb = RGBColor(0, 0, 0)
    p2 = tf.add_paragraph()
    p2.text = body
    p2.font.size = Pt(18)
    p2.font.name = font_name
    p2.space_before = Pt(20)
    p2.font.color.rgb = RGBColor(80, 80, 80)
    return slide

def add_three_block_slide(prs, image_urls, texts, title="", font_name="Arial"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_width = prs.slide_width.inches
    slide_height = prs.slide_height.inches
    block_width = slide_width / 3
    image_height = 2
    text_height = 2
    top_margin = 1.5
    if title:
        title_box = slide.shapes.add_textbox(
            Inches(0.5),
            Inches(0.2),
            Inches(slide_width - 1),
            Inches(1)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.name = font_name
        p.font.color.rgb = RGBColor(0, 102, 204)
    for i in range(3):
        left = Inches(i * block_width)
        img_url = image_urls[i] if i < len(image_urls) else None
        if not img_url or not isinstance(img_url, str) or not img_url.startswith("http"):
            img_url = "https://via.placeholder.com/300x200?text=No+Image"
        img_stream = BytesIO(requests.get(img_url, verify=False).content)
        slide.shapes.add_picture(
            img_stream,
            Inches(i * block_width + (block_width - 2)/2),
            Inches(top_margin),
            Inches(2),
            Inches(image_height)
        )
        textbox = slide.shapes.add_textbox(
            Inches(i * block_width + 0.1),
            Inches(top_margin + image_height + 0.3),
            Inches(block_width - 0.2),
            Inches(text_height)
        )
        tf = textbox.text_frame
        tf.word_wrap = True
        tf.margin_top = 0
        tf.margin_bottom = 0
        p = tf.add_paragraph()
        p.text = texts[i] if i < len(texts) else ''
        p.level = 0
        p.font.size = Pt(16)
        p.font.name = font_name
    return slide

def add_four_block_layout(prs, title_text, blocks, font_name="Arial"):
    assert len(blocks) == 4, "Must provide exactly 4 blocks."
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_width = prs.slide_width.inches
    margin = 0.3
    spacing = 0.2
    title_box = slide.shapes.add_textbox(
        Inches(margin),
        Inches(0.2),
        Inches(slide_width - 2 * margin),
        Inches(1)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.name = font_name
    p.font.color.rgb = RGBColor(0, 102, 204)
    card_width = (slide_width - 2 * margin - 3 * spacing) / 4
    image_height = 1.5
    text_height = 2.5
    top_img = 1.2
    top_txt = top_img + image_height + 0.1
    for i, block in enumerate(blocks):
        left = margin + i * (card_width + spacing)
        img_url = block["image_url"]
        if not img_url or not isinstance(img_url, str) or not img_url.startswith("http"):
            img_url = "https://via.placeholder.com/300x200?text=No+Image"
        img_data = BytesIO(requests.get(img_url, verify=False).content)
        slide.shapes.add_picture(
            img_data,
            Inches(left),
            Inches(top_img),
            Inches(card_width),
            Inches(image_height)
        )
        tb = slide.shapes.add_textbox(
            Inches(left),
            Inches(top_txt),
            Inches(card_width),
            Inches(text_height)
        )
        tf = tb.text_frame
        tf.word_wrap = True
        para = tf.add_paragraph()
        para.text = block["text"]
        para.font.size = Pt(18)
        para.level = 0
        para.font.name = font_name
        para.font.color.rgb = RGBColor(0, 0, 0)
    return slide

def add_basic_content_slide(prs, title, content, font_name="Arial"):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()
    for point in content:
        p = tf.add_paragraph()
        p.text = point
        p.level = 0
        p.font.size = Pt(18)
        p.font.name = font_name
    return slide

# ────── Main Presentation Generation ──────
def slugify_filename(prompt):
    import re
    s = prompt.strip().lower()
    s = re.sub(r'[^a-z0-9\s]', '', s)
    s = re.sub(r'\s+', '_', s)
    return s + '.pptx'

def generate_presentation_from_prompt(prompt, output_path=None, min_slides=5, max_slides=10, font_name="Arial"):
    pres_json = get_presentation_content(prompt, min_slides=min_slides, max_slides=max_slides)
    if not pres_json:
        print("Failed to generate presentation content.")
        return
    prs = Presentation()
    for idx, slide in enumerate(pres_json['slides']):
        layout = slide.get('layout', 'BASIC_CONTENT_SLIDE')
        title = slide.get('title', '')
        content = slide.get('content', [])
        image_query = slide.get('image_query', None)
        if layout in ['1ST_SLIDE']:
            slide_obj = prs.slides.add_slide(prs.slide_layouts[0])
            slide_obj.shapes.title.text = title
            if content:
                slide_obj.placeholders[1].text = '\n'.join(content)
        elif layout in ['BASIC_CONTENT_SLIDE']:
            add_basic_content_slide(prs, title, content, font_name=font_name)
        elif layout in ['2/3RD_IMAGE_LEFT', '2/3RD_IMAGE_RIGHT', 'HALF_IMAGE_LEFT', 'HALF_IMAGE_RIGHT']:
            image_pos = 'right' if 'RIGHT' in layout else 'left'
            image_fraction = 2/3 if '2/3' in layout else 0.5
            img_url = get_unsplash_image_url(image_query)
            add_image_text_slide(prs, img_url, title, '\n'.join(content), image_pos=image_pos, image_fraction=image_fraction, font_name=font_name)
        elif layout == '3_BLOCK_LAYOUT':
            img_urls = get_unsplash_image_urls(image_query)
            add_three_block_slide(prs, img_urls, content, title=title, font_name=font_name)
        elif layout == '4_BLOCK_LAYOUT':
            img_urls = get_unsplash_image_urls(image_query)
            blocks = []
            for i in range(4):
                blocks.append({'image_url': img_urls[i], 'text': content[i] if i < len(content) else ''})
            add_four_block_layout(prs, title, blocks, font_name=font_name)
        else:
            add_basic_content_slide(prs, title, content, font_name=font_name)
    if not output_path:
        output_path = slugify_filename(prompt)
    prs.save(output_path)
    print(f'Presentation saved to {output_path}')

def pptx_to_pdf(pptx_path, pdf_path):
    """
    Converts a PPTX file to PDF using PowerPoint COM automation (requires Windows and MS PowerPoint).
    """
    try:
        import win32com.client
        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        powerpoint.Visible = 1
        presentation = powerpoint.Presentations.Open(pptx_path, WithWindow=False)
        presentation.SaveAs(pdf_path, FileFormat=32)  # 32 = PDF
        presentation.Close()
        powerpoint.Quit()
        return True
    except Exception as e:
        print(f"Error converting PPTX to PDF: {e}")
        return False

# Example usage:
# generate_presentation_from_prompt('The Future of Artificial Intelligence', output_path='ai_future.pptx') 